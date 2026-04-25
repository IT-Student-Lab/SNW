from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn

prs = Presentation('quickscan_43479205d18e.pptx')

def emu_to_in(emu):
    return emu / 914400

for idx, slide in enumerate(prs.slides):
    print(f'=== SLIDE {idx+1} ===')
    for shape in slide.shapes:
        stype = shape.shape_type
        name = shape.name
        l = emu_to_in(shape.left)
        t = emu_to_in(shape.top)
        w = emu_to_in(shape.width)
        h = emu_to_in(shape.height)
        kind = 'shape'
        if stype == 13:
            kind = 'IMAGE'
        elif shape.has_text_frame:
            kind = 'TEXT'
        elif shape.has_table:
            kind = 'TABLE'
        else:
            kind = 'RECT'
        fill_color = ''
        try:
            if hasattr(shape, 'fill') and shape.fill.type is not None:
                fill_color = f' fill=#{shape.fill.fore_color.rgb}'
        except:
            pass
        text_summary = ''
        if shape.has_text_frame:
            all_text = shape.text_frame.text.strip()
            if all_text:
                text_summary = f' "{all_text[:80]}"'
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        f = run.font
                        sz = f.size / 12700 if f.size else None
                        text_summary += f' [font={f.name} size={sz}pt bold={f.bold}'
                        if f.color and f.color.rgb:
                            text_summary += f' color=#{f.color.rgb}'
                        text_summary += ']'
                        break
                if text_summary and '[font' in text_summary:
                    break
        img_info = ''
        if stype == 13 and hasattr(shape, 'image'):
            img_info = f' content_type={shape.image.content_type}'
        print(f'  {kind}: left={l:.2f} top={t:.2f} w={w:.2f} h={h:.2f}{fill_color}{text_summary}{img_info}')
        if shape.has_table:
            table = shape.table
            print(f'    Cols: {len(table.columns)} widths: {[emu_to_in(c.width) for c in table.columns]}')
            print(f'    Rows: {len(table.rows)}')
            for ri, row in enumerate(table.rows):
                for ci, cell in enumerate(row.cells):
                    t2 = cell.text.strip()
                    if t2:
                        fill = ''
                        try:
                            if cell.fill.type is not None:
                                fill = f' bg=#{cell.fill.fore_color.rgb}'
                        except:
                            pass
                        font_info = ''
                        for p in cell.text_frame.paragraphs:
                            for r in p.runs:
                                if r.text.strip():
                                    fo = r.font
                                    sz = fo.size / 12700 if fo.size else None
                                    clr = fo.color.rgb if fo.color and fo.color.rgb else '?'
                                    font_info = f' [size={sz}pt bold={fo.bold} color=#{clr}]'
                                    break
                            if font_info:
                                break
                        print(f'      [{ri},{ci}]: "{t2[:50]}"{fill}{font_info}')
    print()
