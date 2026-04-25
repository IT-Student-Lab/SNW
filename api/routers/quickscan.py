# -*- coding: utf-8 -*-
"""Quickscan router — AI-powered location analysis + PPTX export."""

from __future__ import annotations

import io
import json
from pathlib import Path

import requests as http_requests
from fastapi import APIRouter, Depends, HTTPException, Query
from fastapi.responses import StreamingResponse

from api.deps import get_current_user, get_download_user
from app.config import settings
from app.core import address_to_rd_full, bbox_around_point
from app.core.quickscan import run_quickscan

router = APIRouter(prefix="/api", tags=["quickscan"])


@router.post("/quickscan/{job_id}")
async def quickscan(
    job_id: str,
    address: str | None = Query(None),
    x: float | None = Query(None),
    y: float | None = Query(None),
    radius: float = Query(250),
    _user: str = Depends(get_current_user),
):
    """Run AI quickscan analysis on the generated output of a job.

    Pass either `address` or `x`+`y` so we can build location context.
    """
    # Prevent path traversal
    if "/" in job_id or "\\" in job_id or ".." in job_id:
        raise HTTPException(400, "Ongeldige job id")
    out_dir = Path(settings.output_dir) / job_id
    if not out_dir.is_dir():
        raise HTTPException(404, "Job niet gevonden")

    session = http_requests.Session()
    loc_info = None
    bbox = None

    try:
        if address and len(address.strip()) >= 3:
            loc_info = address_to_rd_full(address.strip(), session=session)
            cx, cy = loc_info["x"], loc_info["y"]
        elif x is not None and y is not None:
            cx, cy = x, y
        else:
            raise HTTPException(400, "Geef een adres of x+y coördinaten op.")

        bbox = tuple(bbox_around_point(cx, cy, radius))
    except HTTPException:
        raise
    except Exception:
        pass  # quickscan can still run without bbox

    display_address = address or (f"RD ({x:.0f}, {y:.0f})" if x and y else "onbekend")

    try:
        sections = run_quickscan(
            out_dir=out_dir,
            loc_info=loc_info,
            adres=display_address,
            radius=radius,
            bbox=bbox,
            session=session,
        )
        return {"job_id": job_id, "sections": sections}
    except Exception as e:
        raise HTTPException(500, f"Quickscan mislukt: {e}")


def _build_quickscan_data(job_id: str, address: str | None, x: float | None,
                          y: float | None, radius: float):
    """Shared helper: resolve location + run quickscan. Returns (sections, out_dir)."""
    if "/" in job_id or "\\" in job_id or ".." in job_id:
        raise HTTPException(400, "Ongeldige job id")
    out_dir = Path(settings.output_dir) / job_id
    if not out_dir.is_dir():
        raise HTTPException(404, "Job niet gevonden")

    session = http_requests.Session()
    loc_info = None
    bbox = None

    try:
        if address and len(address.strip()) >= 3:
            loc_info = address_to_rd_full(address.strip(), session=session)
            cx, cy = loc_info["x"], loc_info["y"]
        elif x is not None and y is not None:
            cx, cy = x, y
        else:
            raise HTTPException(400, "Geef een adres of x+y coördinaten op.")
        bbox = tuple(bbox_around_point(cx, cy, radius))
    except HTTPException:
        raise
    except Exception:
        pass

    display_address = address or (f"RD ({x:.0f}, {y:.0f})" if x and y else "onbekend")

    sections = run_quickscan(
        out_dir=out_dir,
        loc_info=loc_info,
        adres=display_address,
        radius=radius,
        bbox=bbox,
        session=session,
    )
    return sections, out_dir


def _sections_to_pptx(sections: list[dict], out_dir: Path) -> io.BytesIO:
    """Build a PPTX matching the Studio Nico Wissing reference layout exactly."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor

    # Brand colours (exact from reference)
    CLR_PEACH = RGBColor(0xEC, 0xC9, 0xA0)
    CLR_GREEN = RGBColor(0x6B, 0x9B, 0x37)
    CLR_BLUE = RGBColor(0xA8, 0xC8, 0xDE)
    CLR_DARK = RGBColor(0x33, 0x33, 0x33)
    CLR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    CLR_LIGHT_BG = RGBColor(0xF7, 0xF3, 0xED)
    CLR_CAPTION = RGBColor(0x66, 0x66, 0x66)

    SLIDE_W = Inches(13.333)
    SLIDE_H = Inches(7.5)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    def _rect(slide, left, top, width, height, color):
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    def _bg(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _font(run, size=12, bold=False, color=CLR_DARK, name="Calibri"):
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = name

    def _add_text_paragraphs(text_frame, text, size=11, color=CLR_DARK, linkify=False):
        """Split text on newlines and add as paragraphs.

        When *linkify* is True, markdown links [text](url), bare URLs,
        and **bold** are rendered appropriately in the PPTX.
        """
        import re

        # Matches: [text](url), **bold**, or bare URL
        _MD_TOKEN = re.compile(
            r'(\[[^\]]*\]\(https?://[^)]+\)|\*\*[^*]+\*\*|https?://[^\s)\]>"]+)'
        )
        _MD_LINK = re.compile(r'^\[([^\]]*)\]\((https?://[^)]+)\)$')
        _MD_BOLD = re.compile(r'^\*\*(.+)\*\*$')
        _BARE_URL = re.compile(r'^https?://')
        _HEADING = re.compile(r'^#{1,4}\s+')

        text_frame.word_wrap = True
        for i, raw_line in enumerate(text.split("\n")):
            line = _HEADING.sub("", raw_line)  # strip markdown headings
            para = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            if not line.strip():
                continue
            if linkify and _MD_TOKEN.search(line):
                parts = _MD_TOKEN.split(line)
                for part in parts:
                    if not part:
                        continue
                    # Markdown link [text](url)
                    m_link = _MD_LINK.match(part)
                    if m_link:
                        run = para.add_run()
                        run.text = m_link.group(1)
                        _font(run, size=size, color=RGBColor(0x1A, 0x73, 0xE8))
                        run.font.underline = True
                        run.hyperlink.address = m_link.group(2)
                        continue
                    # Bold **text**
                    m_bold = _MD_BOLD.match(part)
                    if m_bold:
                        run = para.add_run()
                        run.text = m_bold.group(1)
                        _font(run, size=size, bold=True, color=color)
                        continue
                    # Bare URL
                    if _BARE_URL.match(part):
                        run = para.add_run()
                        run.text = part
                        _font(run, size=size, color=RGBColor(0x1A, 0x73, 0xE8))
                        run.font.underline = True
                        run.hyperlink.address = part
                        continue
                    # Plain text
                    run = para.add_run()
                    run.text = part
                    _font(run, size=size, color=color)
            else:
                run = para.add_run()
                run.text = line
                _font(run, size=size, color=color)
            para.space_after = Pt(4)

    # ─── SLIDE 1: Title ───────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CLR_WHITE)

    # Three colour bars at top
    for i, clr in enumerate([CLR_PEACH, CLR_GREEN, CLR_BLUE]):
        _rect(slide, Inches(2.7 + i * 2.7), Inches(1.5), Inches(2.5), Inches(0.35), clr)

    # Title + subtitle
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.3), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "AI Quickscan Analyse"
    _font(run, size=40, bold=True, color=CLR_DARK)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "Studio Nico Wissing"
    _font(run2, size=20, color=CLR_GREEN)

    # Green divider
    _rect(slide, Inches(4.0), Inches(4.2), Inches(5.3), Inches(0.04), CLR_GREEN)

    for section in sections:
        title = section.get("title", "")

        # ─── SLIDE 2: Locatie-informatie ──────────────────────────
        if title == "_location_info":
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _bg(slide, CLR_WHITE)

            # Green sidebar
            _rect(slide, Inches(0), Inches(0), Inches(0.25), SLIDE_H, CLR_GREEN)

            # Title
            txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
            run = txBox.text_frame.paragraphs[0].add_run()
            run.text = "Locatie-informatie"
            _font(run, size=28, bold=True, color=CLR_DARK)

            # Peach accent line
            _rect(slide, Inches(0.6), Inches(1.15), Inches(5), Inches(0.03), CLR_PEACH)

            # Info table
            info_items = [
                ("Adres", section.get("display_name", "—")),
                ("Gemeente", section.get("gemeente", "—")),
                ("Provincie", section.get("provincie", "—")),
                ("Woonplaats", section.get("woonplaats", "—")),
                ("Waterschap", section.get("waterschap", "—")),
                ("Buurt", section.get("buurt", "—")),
                ("Straal", f"{section.get('radius', '—')} m"),
            ]

            tbl = slide.shapes.add_table(
                len(info_items), 2, Inches(0.6), Inches(1.5),
                Inches(7), Inches(0.45 * len(info_items))
            ).table
            tbl.columns[0].width = Inches(2.2)
            tbl.columns[1].width = Inches(4.8)

            for row_idx, (label, value) in enumerate(info_items):
                row = tbl.rows[row_idx]
                row.height = Inches(0.45)
                cell_l = row.cells[0]
                cell_v = row.cells[1]
                cell_l.text = ""
                cell_v.text = ""

                run_l = cell_l.text_frame.paragraphs[0].add_run()
                run_l.text = label
                _font(run_l, size=13, bold=True, color=CLR_DARK)
                cell_l.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

                run_v = cell_v.text_frame.paragraphs[0].add_run()
                run_v.text = str(value)
                _font(run_v, size=13, color=CLR_DARK)

                bg_clr = CLR_LIGHT_BG if row_idx % 2 == 0 else CLR_WHITE
                for cell in [cell_l, cell_v]:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = bg_clr
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            continue

        # ─── Normal section slide ─────────────────────────────────
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _bg(slide, CLR_WHITE)

        # Blue sidebar
        _rect(slide, Inches(0), Inches(0), Inches(0.25), SLIDE_H, CLR_BLUE)

        # Section title
        txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.7))
        run = txBox.text_frame.paragraphs[0].add_run()
        run.text = title
        _font(run, size=26, bold=True, color=CLR_DARK)

        # Peach accent line
        _rect(slide, Inches(0.6), Inches(0.95), Inches(4), Inches(0.03), CLR_PEACH)

        # Images layout
        images = section.get("images", [])
        valid_images = [img for img in images if (out_dir / img.get("filename", "")).exists()]
        n_imgs = len(valid_images)

        if n_imgs <= 3:
            # Row of up to 3 images (3.20" each)
            img_top = Inches(1.2)
            IMG_SIZE = Inches(3.2)
            IMG_POSITIONS = [Inches(0.6), Inches(4.7), Inches(8.8)]
            img_count = 0
            for img_info in valid_images:
                img_path = out_dir / img_info.get("filename", "")
                left = IMG_POSITIONS[img_count]
                try:
                    pic = slide.shapes.add_picture(
                        str(img_path), left, img_top, width=IMG_SIZE
                    )
                    if pic.height > IMG_SIZE:
                        ratio = IMG_SIZE / pic.height
                        pic.height = IMG_SIZE
                        pic.width = int(pic.width * ratio)
                    cap = img_info.get("caption", "")
                    if cap:
                        cap_box = slide.shapes.add_textbox(
                            left, img_top + pic.height + Inches(0.05),
                            pic.width, Inches(0.3)
                        )
                        cap_tf = cap_box.text_frame
                        cap_tf.word_wrap = True
                        cap_run = cap_tf.paragraphs[0].add_run()
                        cap_run.text = cap
                        _font(cap_run, size=9, color=CLR_CAPTION)
                        cap_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                    img_count += 1
                except Exception:
                    pass
            text_top = Inches(5.0) if img_count > 0 else Inches(1.2)
        else:
            # 2×2 grid for 4 images (2.85" each)
            IMG_SIZE_4 = Inches(2.85)
            GRID_POSITIONS = [
                (Inches(0.6), Inches(1.2)),
                (Inches(6.6), Inches(1.2)),
                (Inches(0.6), Inches(4.3)),
                (Inches(6.6), Inches(4.3)),
            ]
            img_count = 0
            for idx, img_info in enumerate(valid_images[:4]):
                img_path = out_dir / img_info.get("filename", "")
                left, top = GRID_POSITIONS[idx]
                try:
                    pic = slide.shapes.add_picture(
                        str(img_path), left, top, width=IMG_SIZE_4
                    )
                    if pic.height > IMG_SIZE_4:
                        ratio = IMG_SIZE_4 / pic.height
                        pic.height = IMG_SIZE_4
                        pic.width = int(pic.width * ratio)
                    cap = img_info.get("caption", "")
                    if cap:
                        cap_box = slide.shapes.add_textbox(
                            left, top + pic.height + Inches(0.05),
                            pic.width, Inches(0.3)
                        )
                        cap_tf = cap_box.text_frame
                        cap_tf.word_wrap = True
                        cap_run = cap_tf.paragraphs[0].add_run()
                        cap_run.text = cap
                        _font(cap_run, size=9, color=CLR_CAPTION)
                        cap_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                    img_count += 1
                except Exception:
                    pass
            # No room for text on the image slide; analysis goes on next slide
            text_top = None

        # Analysis text
        analysis = section.get("analysis", "")
        if analysis:
            if text_top is not None:
                text_height = SLIDE_H - text_top - Inches(0.4)
                txBox2 = slide.shapes.add_textbox(
                    Inches(0.6), text_top, Inches(12), text_height
                )
                _add_text_paragraphs(txBox2.text_frame, analysis, linkify=True)
            else:
                # Overflow: put analysis on a continuation slide
                slide_cont = prs.slides.add_slide(prs.slide_layouts[6])
                _bg(slide_cont, CLR_WHITE)
                _rect(slide_cont, Inches(0), Inches(0), Inches(0.25), SLIDE_H, CLR_BLUE)
                txBox_t = slide_cont.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.7))
                run_t = txBox_t.text_frame.paragraphs[0].add_run()
                run_t.text = f"{title} — Analyse"
                _font(run_t, size=24, bold=True, color=CLR_DARK)
                _rect(slide_cont, Inches(0.6), Inches(0.95), Inches(4), Inches(0.03), CLR_PEACH)
                txBox2 = slide_cont.shapes.add_textbox(
                    Inches(0.6), Inches(1.2), Inches(12), Inches(5.8)
                )
                _add_text_paragraphs(txBox2.text_frame, analysis, linkify=True)

        # ─── Omgevingsvisie sub-slide (green sidebar) ─────────────
        omgevingsvisie = section.get("omgevingsvisie")
        if omgevingsvisie:
            slide2 = prs.slides.add_slide(prs.slide_layouts[6])
            _bg(slide2, CLR_WHITE)
            _rect(slide2, Inches(0), Inches(0), Inches(0.25), SLIDE_H, CLR_GREEN)

            txBox = slide2.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.7))
            run = txBox.text_frame.paragraphs[0].add_run()
            run.text = f"{title} \u2014 Omgevingsvisie"
            _font(run, size=24, bold=True, color=CLR_DARK)

            _rect(slide2, Inches(0.6), Inches(0.95), Inches(4), Inches(0.03), CLR_PEACH)

            txBox2 = slide2.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12), Inches(5.8))
            _add_text_paragraphs(txBox2.text_frame, omgevingsvisie, linkify=True)

        # ─── History sub-slide (peach sidebar) ────────────────────
        history_web = section.get("history_web")
        if history_web:
            slide2 = prs.slides.add_slide(prs.slide_layouts[6])
            _bg(slide2, CLR_WHITE)
            _rect(slide2, Inches(0), Inches(0), Inches(0.25), SLIDE_H, CLR_PEACH)

            txBox = slide2.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.7))
            run = txBox.text_frame.paragraphs[0].add_run()
            run.text = f"{title} \u2014 Historische analyse"
            _font(run, size=24, bold=True, color=CLR_DARK)

            _rect(slide2, Inches(0.6), Inches(0.95), Inches(4), Inches(0.03), CLR_PEACH)

            txBox2 = slide2.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12), Inches(5.8))
            _add_text_paragraphs(txBox2.text_frame, history_web, linkify=True)

        # ─── Wikimedia images sub-slide ───────────────────────────
        wiki_images = section.get("wikimedia_images") or []
        displayable = [w for w in wiki_images if (out_dir / w.get("local_file", "")).exists()]
        if displayable:
            slide3 = prs.slides.add_slide(prs.slide_layouts[6])
            _bg(slide3, CLR_WHITE)
            _rect(slide3, Inches(0), Inches(0), Inches(0.25), SLIDE_H, CLR_PEACH)

            txBox = slide3.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.7))
            run = txBox.text_frame.paragraphs[0].add_run()
            run.text = f"{title} — Wikimedia Commons"
            _font(run, size=24, bold=True, color=CLR_DARK)

            _rect(slide3, Inches(0.6), Inches(0.95), Inches(4), Inches(0.03), CLR_PEACH)

            # Grid: up to 6 images in 3×2
            WIKI_SIZE = Inches(3.5)
            WIKI_COLS = [Inches(0.6), Inches(4.7), Inches(8.8)]
            WIKI_ROWS = [Inches(1.2), Inches(4.3)]
            for wi, wimg in enumerate(displayable[:6]):
                col = wi % 3
                row = wi // 3
                wleft = WIKI_COLS[col]
                wtop = WIKI_ROWS[row]
                wpath = out_dir / wimg["local_file"]
                try:
                    wpic = slide3.shapes.add_picture(
                        str(wpath), wleft, wtop, width=WIKI_SIZE
                    )
                    if wpic.height > Inches(2.8):
                        ratio = Inches(2.8) / wpic.height
                        wpic.height = Inches(2.8)
                        wpic.width = int(wpic.width * ratio)
                    wtitle = wimg.get("title", "")
                    if wtitle:
                        wcap = slide3.shapes.add_textbox(
                            wleft, wtop + wpic.height + Inches(0.05),
                            wpic.width, Inches(0.3)
                        )
                        wcap_tf = wcap.text_frame
                        wcap_tf.word_wrap = True
                        wcap_run = wcap_tf.paragraphs[0].add_run()
                        wcap_run.text = wtitle[:60]
                        _font(wcap_run, size=8, color=CLR_CAPTION)
                        wcap_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                except Exception:
                    pass

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


@router.get("/quickscan/{job_id}/cached")
async def quickscan_cached(
    job_id: str,
    _user: str = Depends(get_current_user),
):
    """Return cached quickscan results if available."""
    if "/" in job_id or "\\" in job_id or ".." in job_id:
        raise HTTPException(400, "Ongeldige job id")
    cache_file = Path(settings.output_dir) / job_id / "quickscan.json"
    if not cache_file.exists():
        return {"sections": None}
    sections = json.loads(cache_file.read_text(encoding="utf-8"))
    return {"sections": sections}


@router.get("/quickscan/{job_id}/export")
async def quickscan_export_pptx(
    job_id: str,
    address: str | None = Query(None),
    x: float | None = Query(None),
    y: float | None = Query(None),
    radius: float = Query(250),
    _user: str = Depends(get_download_user),
):
    """Export quickscan results as a PowerPoint (.pptx) file.

    Uses cached quickscan.json from disk, or falls back to running
    the full quickscan.
    """
    if "/" in job_id or "\\" in job_id or ".." in job_id:
        raise HTTPException(400, "Ongeldige job id")
    out_dir = Path(settings.output_dir) / job_id
    if not out_dir.is_dir():
        raise HTTPException(404, "Job niet gevonden")

    # Try loading cached results from disk first
    cache_file = out_dir / "quickscan.json"
    if cache_file.exists():
        sections = json.loads(cache_file.read_text(encoding="utf-8"))
    else:
        try:
            sections, out_dir = _build_quickscan_data(job_id, address, x, y, radius)
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(500, f"Quickscan mislukt: {e}")

    buf = _sections_to_pptx(sections, out_dir)

    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": f"attachment; filename=quickscan_{job_id}.pptx"
        },
    )
